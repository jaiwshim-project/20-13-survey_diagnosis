"""
대한치과위생사협회 조직진단 결과 리포트 → PowerPoint 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math, os

# ── 색상 팔레트 ──────────────────────────────
C_NAVY    = RGBColor(0x0f, 0x27, 0x44)   # primary
C_BLUE    = RGBColor(0x25, 0x63, 0xeb)   # accent
C_BLUE2   = RGBColor(0x1d, 0x4e, 0xd8)
C_RED     = RGBColor(0xdc, 0x26, 0x26)
C_ORANGE  = RGBColor(0xd9, 0x77, 0x06)
C_AMBER   = RGBColor(0xb4, 0x53, 0x09)
C_GREEN   = RGBColor(0x05, 0x96, 0x69)
C_GOLD    = RGBColor(0xd4, 0xa8, 0x53)
C_WHITE   = RGBColor(0xff, 0xff, 0xff)
C_GRAY50  = RGBColor(0xf8, 0xfa, 0xfc)
C_GRAY100 = RGBColor(0xf1, 0xf5, 0xf9)
C_GRAY200 = RGBColor(0xe2, 0xe8, 0xf0)
C_GRAY500 = RGBColor(0x64, 0x74, 0x8b)
C_DARK    = RGBColor(0x1e, 0x29, 0x3b)

W = Inches(13.33)   # 와이드 슬라이드 너비
H = Inches(7.5)     # 높이

LOGO_PATH = os.path.join(os.path.dirname(__file__), "images", "kdha-logo.png")

# ── 헬퍼 함수 ────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=Pt(14), bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, v_anchor=None, wrap=True,
             font_name="맑은 고딕"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb

def add_rect_text(slide, text, x, y, w, h,
                  fill=C_BLUE, text_color=C_WHITE,
                  font_size=Pt(13), bold=False,
                  align=PP_ALIGN.CENTER, radius=False):
    """색상 박스 + 텍스트 합성"""
    r = add_rect(slide, x, y, w, h, fill=fill)
    tf = r.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.color.rgb = text_color
    run.font.name      = "맑은 고딕"
    return r

def score_color(score):
    if score >= 70: return C_GREEN
    if score >= 60: return C_AMBER
    if score >= 50: return C_ORANGE
    return C_RED

def grade_text(score):
    if score >= 70: return "양호"
    if score >= 60: return "보통"
    if score >= 50: return "미흡"
    if score >= 45: return "취약"
    return "위험"

def grade_color(score):
    return score_color(score)

# ── 슬라이드 배경 채우기 ─────────────────────
def bg(slide, color):
    add_rect(slide, 0, 0, W, H, fill=color)

def nav_bar(slide, title=""):
    """상단 네이비 바"""
    add_rect(slide, 0, 0, W, Inches(0.55), fill=C_NAVY)
    if title:
        add_text(slide, title,
                 Inches(0.4), Inches(0.05), Inches(10), Inches(0.45),
                 font_size=Pt(16), bold=True, color=C_WHITE)

def page_num(slide, n, total):
    add_text(slide, f"{n} / {total}",
             Inches(12.4), Inches(7.1), Inches(0.9), Inches(0.3),
             font_size=Pt(9), color=C_GRAY500, align=PP_ALIGN.RIGHT)

# ── 레이더 차트 (도형으로 근사) ──────────────
def draw_radar(slide, cx, cy, max_r, scores):
    """7각형 레이더를 pptx 도형으로 그린다"""
    n = len(scores)
    def pt(i, r):
        a = -math.pi/2 + i * 2*math.pi/n
        return cx + r*math.cos(a), cy + r*math.sin(a)

    # 그리드 링
    for level in [0.25, 0.5, 0.75, 1.0]:
        pts = [pt(i, max_r*level) for i in range(n)]
        from pptx.util import Pt as PPt
        sp = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.01), Inches(0.01))
        sp.fill.background()
        # 다각형 직접 그리기는 add_connector 한계 → connector 라인으로 대신
        for i in range(n):
            x1, y1 = pts[i]
            x2, y2 = pts[(i+1) % n]
            conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
            conn.line.color.rgb = C_GRAY200 if level < 1.0 else RGBColor(0xcb,0xd5,0xe1)
            conn.line.width = Pt(0.75) if level < 1.0 else Pt(1.2)

    # 축 선
    for i in range(n):
        x1, y1 = cx, cy
        x2, y2 = pt(i, max_r)
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = C_GRAY200
        conn.line.width = Pt(0.5)

    # 데이터 다각형 (연결선)
    data_pts = [pt(i, max_r * scores[i]/100) for i in range(n)]
    for i in range(n):
        x1, y1 = data_pts[i]
        x2, y2 = data_pts[(i+1) % n]
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = C_BLUE
        conn.line.width = Pt(2.2)

    # 데이터 포인트 원
    r_dot = Inches(0.08)
    for i, (dx, dy) in enumerate(data_pts):
        circ = slide.shapes.add_shape(9, dx-r_dot/2, dy-r_dot/2, r_dot, r_dot)
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_BLUE
        circ.line.color.rgb = C_WHITE
        circ.line.width = Pt(1.5)

# ── 데이터 ───────────────────────────────────
AREAS = [
    {"no": 1, "name": "공공 미션 작동력",       "score": 62,
     "strength": "협회 공공 미션 구성원 인식률 73%.\n주요 사업 기획 시 공공성 원칙 표방 문화 존재.",
     "weakness": "의사결정 시 공공성 기준 적용률 51%.\n회원 이익·공공 이익 충돌 시 판단 기준 불명확 67%.",
     "quote": "미션은 다 알지만, 안건 결정 시 공공 미션 기준으로 토론하는 건 잘 못 봤어요."},
    {"no": 2, "name": "전문직 자율성 기준",     "score": 54,
     "strength": "전문직 자율성 가치 공감률 79%.\n전문 판단 존중 문화 부분 형성.",
     "weakness": "자율 판단·승인 영역 경계 명확성 38%.\n명문화된 위임 기준서 보유율 27%.",
     "quote": "어디까지 내가 결정할 수 있는지 몰라요. 매번 눈치를 봐요."},
    {"no": 3, "name": "선출 거버넌스 연속성",   "score": 48,
     "strength": "신규 집행부의 업무 파악 의지·학습 동기 높음.\n비공식 소통 시도 사례 일부 확인.",
     "weakness": "인수인계 문서·절차 보유율 23%.\n집행부 교체 시 4~6개월 공백 반복. 중장기 정책 연속성 구조 전무.",
     "quote": "새 집행부가 오면 처음 6개월은 파악하는 데 써요. 그동안 중요한 결정이 미뤄지죠."},
    {"no": 4, "name": "사무국-임원 역할 경계",  "score": 71,
     "strength": "역할 경계 인식 일치율 68% (7개 영역 최고).\n일상 업무에서 역할 구분 비교적 잘 이루어짐.",
     "weakness": "긴급 상황 역할 혼란 경험 44%.\n역할 경계 문서 존재하나 일관성 부족 51%.",
     "quote": "급하거나 중요한 일이 생기면 임원이 개입하게 돼요. 담당자는 어떻게 해야 할지 모르게 되죠."},
    {"no": 5, "name": "세대 간 일의 기준",      "score": 45,
     "strength": "세대 간 갈등 존재 양측 인식, 문제 공유 기반 형성.\n대화 의지 구성원 비율 64%.",
     "weakness": "세대 초월 공통 업무 기준 보유율 31%.\n20~30대 67% '기준이 사람마다 다르다' 응답.",
     "quote": "선배들은 태도를 보고, 저는 결과로 평가받고 싶은데 기준이 어디에도 없어요."},
    {"no": 6, "name": "조직 기억 유지 구조",    "score": 39,
     "strength": "베테랑 구성원 개인 아카이빙 노력으로 일부 보존.\n기록 필요성 공감률 82%.",
     "weakness": "지식 관리 시스템 보유율 18% (7개 영역 최저).\n온보딩 기간 평균 6개월 이상. 성공·실패 정리 전무.",
     "quote": "담당자가 바뀌면 처음부터 다시 시작이에요. 같은 실수를 반복하게 됩니다."},
    {"no": 7, "name": "신뢰/공정성 체감 구조",  "score": 58,
     "strength": "의사결정 결과 공유율 71%.\n구성원 의견 수렴 채널 형식적 운영.",
     "weakness": "의사결정 과정 투명성 체감률 47%.\n'이유 없이 통보받는다' 응답 59%. 반영 피드백 채널 부재.",
     "quote": "결론은 알려줘요. 왜 그렇게 결정됐는지 모르고, 내 의견 반영 여부도 피드백이 없어요."},
]

WORKSHOPS = [
    {"no":"WS-01","color":C_RED,   "target":"조직 기억 유지 구조 39점",
     "title":"우리 협회의 기억을 살린다",
     "sub":"조직 기억 복원 & 지식 관리 체계 구축",
     "timing":"즉시 착수 (1개월 내)", "days":"2일",
     "day1":"핵심 담당자 집단 인터뷰 · 업무 노하우·실패 경험 추출",
     "day2":"지식 카테고리 분류 · 온보딩 매뉴얼 초안 공동 작성",
     "output":"협회 핵심 업무 지식 아카이브 · 온보딩 가이드북 v1.0 · 지식 관리 체계 설계안"},
    {"no":"WS-02","color":C_RED,   "target":"세대 간 일의 기준 45점",
     "title":"우리는 어떻게 일할 것인가",
     "sub":"세대 통합형 공통 업무 기준 정립",
     "timing":"즉시 착수 (2개월 내)", "days":"2일 + 사후 검토",
     "day1":"세대별 일의 가치관 카드 워크숍 · 갈등 지점 시각화",
     "day2":"공통 업무 기준 항목 도출 · 평가 기준 공동 합의",
     "output":"세대 통합형 업무 기준서 · 공통 평가 기준 매트릭스 · 갈등 조정 프로세스"},
    {"no":"WS-03","color":C_ORANGE,"target":"선출 거버넌스 연속성 48점",
     "title":"바뀌어도 흔들리지 않는다",
     "sub":"집행부 연속성 체계 설계",
     "timing":"단기 (3개월 내)", "days":"2일",
     "day1":"집행부 교체 시나리오 시뮬레이션 · 공백 리스크 도출",
     "day2":"인수인계 표준 체계 설계 · 중장기 정책 연속성 로드맵",
     "output":"집행부 인수인계 표준 매뉴얼 · 섀도 캐비닛 운영 규정 초안 · 중장기 정책 승계 지도"},
    {"no":"WS-04","color":C_ORANGE,"target":"전문직 자율성 기준 54점",
     "title":"내 판단의 경계를 안다",
     "sub":"전문직 자율성 & 위임 기준 명문화",
     "timing":"단기 (4개월 내)", "days":"1.5일",
     "day1":"실제 갈등 사례 분석 · 자율·통제 스펙트럼 매핑",
     "day2":"업무 유형별 위임 권한 테이블 작성 · 책임 기준 합의",
     "output":"업무 위임 권한 테이블 · 자율 판단 기준서 · 자율성-책임 매핑 차트"},
    {"no":"WS-05","color":C_BLUE,  "target":"공공 미션 작동력 62점",
     "title":"미션이 실제로 작동하게 한다",
     "sub":"공공 미션 기반 의사결정 체계",
     "timing":"단기 (5개월 내)", "days":"1일",
     "day1":"과거 안건으로 공공 미션 적용 여부 재분석",
     "day2":"공공 미션 의사결정 5단계 가이드 공동 작성",
     "output":"공공 미션 의사결정 가이드 · 미션 기반 사업 평가 체크리스트 · 회의 운영 원칙"},
    {"no":"WS-06","color":C_BLUE,  "target":"신뢰/공정성 체감 구조 58점",
     "title":"투명하게, 공정하게",
     "sub":"의사결정 투명성 & 신뢰 구조 재설계",
     "timing":"단기 (6개월 내)", "days":"1일",
     "day1":"불신 발생 구조 분석 · 투명성 저해 요인 도출",
     "day2":"의사결정 공개 기준 설계 · 구성원 의견 반영 채널 설계",
     "output":"의사결정 투명성 운영 원칙 · 구성원 의견 반영 채널 설계안 · 공정성 체감 지표"},
]

TOTAL = 20   # 총 슬라이드 수 (대략)

# ══════════════════════════════════════════════
# 슬라이드 생성
# ══════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # 완전 빈 레이아웃

sn = [0]   # 슬라이드 번호 카운터

def new_slide():
    sn[0] += 1
    return prs.slides.add_slide(blank)

# ────────────────────────────────────────────
# 슬라이드 1 : 표지
# ────────────────────────────────────────────
sl = new_slide()
# 배경 그라디언트 효과 (단색 근사)
add_rect(sl, 0, 0, W, H, fill=C_NAVY)
add_rect(sl, 0, 0, W*0.55, H, fill=RGBColor(0x0d,0x1f,0x35))

# 왼쪽 강조 바
add_rect(sl, Inches(0.55), Inches(1.5), Inches(0.07), Inches(4.5), fill=C_BLUE)

# SAMPLE 배지
add_rect_text(sl, "★  SAMPLE REPORT",
              Inches(0.75), Inches(1.6), Inches(2.4), Inches(0.38),
              fill=RGBColor(0x1a,0x3a,0x5c), text_color=C_GOLD,
              font_size=Pt(11), bold=True)

# 협회 로고
if os.path.exists(LOGO_PATH):
    sl.shapes.add_picture(LOGO_PATH,
                          Inches(0.75), Inches(2.15),
                          height=Inches(0.55))

# 제목
add_text(sl, "대한치과위생사협회",
         Inches(0.75), Inches(2.9), Inches(6.5), Inches(1.0),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_text(sl, "조직진단 결과 리포트",
         Inches(0.75), Inches(3.85), Inches(6.5), Inches(0.7),
         font_size=Pt(24), bold=False, color=RGBColor(0xb0,0xc8,0xf0))

# 구분선
add_rect(sl, Inches(0.75), Inches(4.6), Inches(5.5), Inches(0.03), fill=RGBColor(0x2a,0x4a,0x70))

# 메타 정보
metas = [
    ("진단 기간",   "2025년 11월 4일 ~ 11월 15일"),
    ("참여 인원",   "32명 / 35명  (응답률 91.4%)"),
    ("진단 영역",   "7대 조직 진단 영역"),
    ("발  행  일",  "2025년 11월 25일"),
]
for i, (k, v) in enumerate(metas):
    y = Inches(4.8) + i * Inches(0.42)
    add_text(sl, k, Inches(0.75), y, Inches(1.4), Inches(0.38),
             font_size=Pt(9), color=C_GRAY500)
    add_text(sl, v, Inches(2.1), y, Inches(4.5), Inches(0.38),
             font_size=Pt(11), bold=True, color=C_WHITE)

# 오른쪽 장식 원
add_rect(sl, Inches(8.5), Inches(-1), Inches(5), Inches(5),
         fill=RGBColor(0x17,0x35,0x60))
add_rect(sl, Inches(9.5), Inches(3.5), Inches(4), Inches(4),
         fill=RGBColor(0x12,0x2a,0x50))

# 우측 요약 숫자
stats = [("54점", "종합 진단 점수"), ("7개", "진단 영역"), ("32명", "참여 구성원")]
for i,(num,lab) in enumerate(stats):
    x = Inches(8.2) + i * Inches(1.7)
    add_text(sl, num, x, Inches(2.8), Inches(1.6), Inches(0.8),
             font_size=Pt(28), bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    add_text(sl, lab, x, Inches(3.55), Inches(1.6), Inches(0.4),
             font_size=Pt(9), color=RGBColor(0x90,0xaa,0xcc), align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────
# 슬라이드 2 : 목차
# ────────────────────────────────────────────
sl = new_slide()
bg(sl, C_GRAY50)
nav_bar(sl, "목  차")

items = [
    ("01", "종합 진단 요약",       "Executive Summary"),
    ("02", "7대 영역 점수 요약",   "Score Summary"),
    ("03", "레이더 분석 차트",     "Radar Analysis"),
    ("04", "영역별 상세 결과",     "Detailed Results  (Area 1 – 7)"),
    ("05", "우선순위 개선 과제",   "Priority Action Plan"),
    ("06", "워크숍 솔루션 제안",   "Workshop Solutions  (WS-01 – 06)"),
    ("07", "워크숍 추진 일정",     "Workshop Roadmap"),
    ("08", "결론 및 다음 단계",    "Conclusion"),
]

for i, (no, ko, en) in enumerate(items):
    row = i % 4
    col = i // 4
    x = Inches(0.7) + col * Inches(6.4)
    y = Inches(1.1) + row * Inches(1.5)

    add_rect(sl, x, y, Inches(5.8), Inches(1.25),
             fill=C_WHITE, line=C_GRAY200, line_w=Pt(1))
    add_rect(sl, x, y, Inches(0.65), Inches(1.25), fill=C_NAVY)
    add_text(sl, no, x, y, Inches(0.65), Inches(1.25),
             font_size=Pt(15), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, ko, x+Inches(0.75), y+Inches(0.12), Inches(4.9), Inches(0.6),
             font_size=Pt(15), bold=True, color=C_NAVY)
    add_text(sl, en, x+Inches(0.75), y+Inches(0.68), Inches(4.9), Inches(0.45),
             font_size=Pt(10), color=C_GRAY500)

page_num(sl, 2, TOTAL)


# ────────────────────────────────────────────
# 슬라이드 3 : 종합 진단 요약
# ────────────────────────────────────────────
sl = new_slide()
bg(sl, C_WHITE)
nav_bar(sl, "종합 진단 요약  |  Executive Summary")

# 왼쪽 점수 원
cx, cy, r = Inches(2.4), Inches(4.0), Inches(1.45)
add_rect(sl, cx-r, cy-r, r*2, r*2, fill=C_BLUE)
# 흰 내부 원
ri = Inches(1.05)
add_rect(sl, cx-ri, cy-ri, ri*2, ri*2, fill=C_WHITE)
add_text(sl, "54", cx-ri, cy-ri-Inches(0.1), ri*2, ri*2,
         font_size=Pt(46), bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "/ 100점", cx-ri, cy+Inches(0.55), ri*2, Inches(0.4),
         font_size=Pt(11), color=C_GRAY500, align=PP_ALIGN.CENTER)

add_rect_text(sl, "보통 (미흡 경계)",
              cx-Inches(1.0), cy+ri+Inches(0.15), Inches(2.0), Inches(0.38),
              fill=RGBColor(0xfe,0xf3,0xc7), text_color=C_ORANGE,
              font_size=Pt(12), bold=True)

add_text(sl, "종합 점수", cx-Inches(1.2), Inches(0.7), Inches(2.4), Inches(0.5),
         font_size=Pt(13), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

# 오른쪽 요약 텍스트
add_text(sl, "진단 핵심 요약",
         Inches(4.5), Inches(0.7), Inches(8.5), Inches(0.5),
         font_size=Pt(14), bold=True, color=C_NAVY)

summary_lines = [
    "대한치과위생사협회는 7대 진단 영역 전반에서 평균 54점을 기록하였습니다.",
    "사무국-임원 역할 경계(71점)가 상대적으로 높은 반면,",
    "조직 기억 유지 구조(39점)·세대 간 일의 기준(45점)에서 심각한 취약점이 확인되었습니다.",
    "선출 거버넌스 연속성(48점)은 집행부 교체마다 반복되는 혼란의 핵심 원인입니다.",
]
for i, line in enumerate(summary_lines):
    add_text(sl, line,
             Inches(4.5), Inches(1.3)+i*Inches(0.44), Inches(8.5), Inches(0.4),
             font_size=Pt(11.5), color=C_DARK)

# 강점/위험 박스
for col, (label, items_list, fill, tc) in enumerate([
    ("▲ 강점 영역", ["사무국-임원 역할 경계  71점", "공공 미션 작동력  62점"],
     RGBColor(0xf0,0xfd,0xf4), C_GREEN),
    ("▼ 위험 영역", ["조직 기억 유지 구조  39점", "세대 간 일의 기준  45점"],
     RGBColor(0xfe,0xf2,0xf2), C_RED),
]):
    x = Inches(4.5) + col * Inches(4.1)
    add_rect(sl, x, Inches(3.2), Inches(3.8), Inches(1.5),
             fill=fill, line=C_GRAY200, line_w=Pt(0.5))
    add_text(sl, label, x+Inches(0.15), Inches(3.3), Inches(3.5), Inches(0.4),
             font_size=Pt(11), bold=True, color=tc)
    for j, itm in enumerate(items_list):
        add_text(sl, "• " + itm,
                 x+Inches(0.15), Inches(3.75)+j*Inches(0.38), Inches(3.5), Inches(0.35),
                 font_size=Pt(10.5), color=C_DARK)

# 7개 영역 가로 점수 바
add_text(sl, "영역별 점수 한눈에 보기",
         Inches(0.5), Inches(5.05), Inches(12), Inches(0.4),
         font_size=Pt(12), bold=True, color=C_NAVY)

bar_w_total = Inches(10.5)
for i, a in enumerate(AREAS):
    y = Inches(5.55) + i * Inches(0.255)
    # 이름
    add_text(sl, f"{a['no']}. {a['name']}",
             Inches(0.5), y, Inches(2.8), Inches(0.24),
             font_size=Pt(8.5), color=C_DARK)
    # 트랙
    add_rect(sl, Inches(3.35), y+Inches(0.05), bar_w_total, Inches(0.14), fill=C_GRAY100)
    # 채움
    add_rect(sl, Inches(3.35), y+Inches(0.05),
             bar_w_total * a['score']/100, Inches(0.14),
             fill=score_color(a['score']))
    # 점수
    add_text(sl, f"{a['score']}점",
             Inches(3.4) + bar_w_total * a['score']/100, y,
             Inches(0.7), Inches(0.24),
             font_size=Pt(8), bold=True, color=score_color(a['score']))

page_num(sl, 3, TOTAL)


# ────────────────────────────────────────────
# 슬라이드 4 : 영역별 점수 요약 테이블
# ────────────────────────────────────────────
sl = new_slide()
bg(sl, C_GRAY50)
nav_bar(sl, "7대 영역 점수 요약  |  Score Summary")

headers = ["No.", "진단 영역", "점수", "등급", "핵심 이슈"]
col_x   = [Inches(0.3), Inches(0.85), Inches(3.75), Inches(4.55), Inches(5.55)]
col_w   = [Inches(0.5), Inches(2.85), Inches(0.75), Inches(0.95), Inches(7.35)]
row_h   = Inches(0.64)

# 헤더
for j, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(sl, x, Inches(0.65), w, Inches(0.5), fill=C_NAVY)
    add_text(sl, hdr, x, Inches(0.65), w, Inches(0.5),
             font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

issues = [
    "미션 인식은 있으나 의사결정 적용 기준 미흡",
    "자율 판단 영역과 승인 영역 경계 불명확",
    "인수인계 체계 부재, 집행부 교체 시 공백 반복",
    "기본 역할 구분 형성, 긴급 상황 시 혼란 존재",
    "세대별 업무 기준 충돌, 공통 기준 미정립",
    "체계적 기록 시스템 전무, 지식 개인화 심각",
    "결과 공유는 됨, 의사결정 과정 투명성 부족",
]

for i, (a, issue) in enumerate(zip(AREAS, issues)):
    y = Inches(1.2) + i * row_h
    row_fill = C_WHITE if i % 2 == 0 else C_GRAY50

    for j, (x, w) in enumerate(zip(col_x, col_w)):
        add_rect(sl, x, y, w, row_h-Inches(0.04),
                 fill=row_fill, line=C_GRAY200, line_w=Pt(0.3))

    sc = score_color(a['score'])
    cells = [str(a['no']), a['name'], f"{a['score']}점",
             grade_text(a['score']), issue]
    aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER,
              PP_ALIGN.CENTER, PP_ALIGN.LEFT]
    bolds  = [True, True, True, True, False]

    for j, (txt, x, w, aln, bld) in enumerate(zip(cells, col_x, col_w, aligns, bolds)):
        clr = sc if j in (2, 3) else C_DARK
        add_text(sl, txt, x+Inches(0.06), y+Inches(0.12), w-Inches(0.1), row_h-Inches(0.18),
                 font_size=Pt(10.5 if j != 4 else 9.5),
                 bold=bld, color=clr, align=aln)

page_num(sl, 4, TOTAL)


# ────────────────────────────────────────────
# 슬라이드 5 : 레이더 차트
# ────────────────────────────────────────────
sl = new_slide()
bg(sl, C_WHITE)
nav_bar(sl, "7대 영역 레이더 분석  |  Radar Analysis")

add_text(sl, "이상적 수준(100점)과 현재 진단 수준의 차이를 시각화한 레이더 차트입니다.",
         Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         font_size=Pt(11), color=C_GRAY500)

cx  = Inches(6.65)
cy  = Inches(4.2)
max_r = Inches(2.7)
scores = [a['score'] for a in AREAS]
draw_radar(sl, cx, cy, max_r, scores)

# 레이블
label_r = max_r + Inches(0.65)
label_names = [a['name'] for a in AREAS]
for i, name in enumerate(label_names):
    ang = -math.pi/2 + i * 2*math.pi/7
    lx = cx + label_r * math.cos(ang)
    ly = cy + label_r * math.sin(ang)
    sc = AREAS[i]['score']
    w_lbl, h_lbl = Inches(1.4), Inches(0.55)
    add_text(sl, name, lx-w_lbl/2, ly-h_lbl/2, w_lbl, h_lbl,
             font_size=Pt(9), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    # 점수 표시
    dr = max_r * sc/100
    dx = cx + dr * math.cos(ang)
    dy = cy + dr * math.sin(ang)
    add_rect_text(sl, f"{sc}점",
                  dx-Inches(0.24), dy-Inches(0.18), Inches(0.48), Inches(0.28),
                  fill=score_color(sc), text_color=C_WHITE,
                  font_size=Pt(8), bold=True)

# 범례
legend_items = [("양호 (70+)", C_GREEN), ("보통 (60~69)", C_AMBER),
                ("미흡 (50~59)", C_ORANGE), ("취약/위험 (~49)", C_RED)]
for i, (lbl, clr) in enumerate(legend_items):
    x = Inches(0.4) + i * Inches(1.85)
    add_rect(sl, x, Inches(7.05), Inches(0.18), Inches(0.18), fill=clr)
    add_text(sl, lbl, x+Inches(0.24), Inches(7.04), Inches(1.55), Inches(0.22),
             font_size=Pt(8.5), color=C_DARK)

page_num(sl, 5, TOTAL)


# ────────────────────────────────────────────
# 슬라이드 6~12 : 영역별 상세 (각 영역 1장)
# ────────────────────────────────────────────
for idx, a in enumerate(AREAS):
    sl = new_slide()
    bg(sl, C_WHITE)

    sc   = a['score']
    sc_c = score_color(sc)
    grd  = grade_text(sc)

    # 상단 컬러 헤더
    add_rect(sl, 0, 0, W, Inches(1.5), fill=sc_c)
    add_text(sl, f"영역 {a['no']}  |  {a['name']}",
             Inches(0.5), Inches(0.12), Inches(9), Inches(0.55),
             font_size=Pt(20), bold=True, color=C_WHITE)
    add_text(sl, f"Diagnosis Area {a['no']} of 7",
             Inches(0.5), Inches(0.68), Inches(9), Inches(0.4),
             font_size=Pt(11), color=RGBColor(0xff,0xff,0xff))

    # 점수 원 (우상단)
    add_rect(sl, Inches(11.0), Inches(0.1), Inches(1.6), Inches(1.3), fill=sc_c)
    add_text(sl, f"{sc}점", Inches(11.0), Inches(0.1), Inches(1.6), Inches(0.85),
             font_size=Pt(34), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, grd, Inches(11.0), Inches(0.9), Inches(1.6), Inches(0.4),
             font_size=Pt(12), color=C_WHITE, align=PP_ALIGN.CENTER)

    # 점수 바
    bar_y = Inches(1.7)
    add_text(sl, "진단 점수", Inches(0.5), bar_y, Inches(1.2), Inches(0.35),
             font_size=Pt(10), color=C_GRAY500)
    add_rect(sl, Inches(1.7), bar_y+Inches(0.07), Inches(10.0), Inches(0.2), fill=C_GRAY100)
    add_rect(sl, Inches(1.7), bar_y+Inches(0.07), Inches(10.0)*sc/100, Inches(0.2), fill=sc_c)
    add_text(sl, f"{sc} / 100",
             Inches(1.7)+Inches(10.0)*sc/100+Inches(0.05), bar_y,
             Inches(1.0), Inches(0.35),
             font_size=Pt(10), bold=True, color=sc_c)

    # 강점 박스
    add_rect(sl, Inches(0.4), Inches(2.2), Inches(6.0), Inches(2.1),
             fill=RGBColor(0xf0,0xfd,0xf4), line=C_GREEN, line_w=Pt(1.5))
    add_rect(sl, Inches(0.4), Inches(2.2), Inches(0.08), Inches(2.1), fill=C_GREEN)
    add_text(sl, "▲  강점",
             Inches(0.6), Inches(2.25), Inches(5.6), Inches(0.38),
             font_size=Pt(11), bold=True, color=C_GREEN)
    add_text(sl, a['strength'],
             Inches(0.6), Inches(2.65), Inches(5.7), Inches(1.5),
             font_size=Pt(10.5), color=C_DARK)

    # 약점 박스
    add_rect(sl, Inches(6.7), Inches(2.2), Inches(6.0), Inches(2.1),
             fill=RGBColor(0xfe,0xf2,0xf2), line=C_RED, line_w=Pt(1.5))
    add_rect(sl, Inches(6.7), Inches(2.2), Inches(0.08), Inches(2.1), fill=C_RED)
    add_text(sl, "▼  약점",
             Inches(6.9), Inches(2.25), Inches(5.6), Inches(0.38),
             font_size=Pt(11), bold=True, color=C_RED)
    add_text(sl, a['weakness'],
             Inches(6.9), Inches(2.65), Inches(5.7), Inches(1.5),
             font_size=Pt(10.5), color=C_DARK)

    # 인용구
    add_rect(sl, Inches(0.4), Inches(4.55), Inches(12.3), Inches(1.15),
             fill=RGBColor(0xef,0xf6,0xff), line=C_BLUE, line_w=Pt(1.5))
    add_rect(sl, Inches(0.4), Inches(4.55), Inches(0.08), Inches(1.15), fill=C_BLUE)
    add_text(sl, '"' + a['quote'] + '"',
             Inches(0.65), Inches(4.6), Inches(11.9), Inches(1.05),
             font_size=Pt(11), color=C_DARK)

    # 개선 방향 힌트
    add_text(sl, "→  워크숍을 통한 개선 가능 영역  |  상세 솔루션은 워크숍 제안 섹션 참조",
             Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.4),
             font_size=Pt(9.5), color=C_BLUE)

    page_num(sl, 6+idx, TOTAL)


# ────────────────────────────────────────────
# 슬라이드 13 : 우선순위 개선 과제
# ────────────────────────────────────────────
sl = new_slide()
bg(sl, C_GRAY50)
nav_bar(sl, "우선순위 개선 과제  |  Priority Action Plan")

add_text(sl, "긴급도·중요도 기준으로 도출한 단계별 개선 과제입니다.",
         Inches(0.5), Inches(0.65), Inches(12), Inches(0.38),
         font_size=Pt(11), color=C_GRAY500)

phases = [
    ("즉시 착수", "1 ~ 3개월", C_RED, [
        "조직 기억 긴급 수집 워크숍 실시 (핵심 담당자 인터뷰)",
        "집행부 인수인계 표준 체계 초안 수립",
        "세대 간 갈등 현황 심층 인터뷰 실시",
        "긴급 상황 역할 대응 매뉴얼 마련",
    ]),
    ("단기 과제", "3 ~ 6개월", C_ORANGE, [
        "전문직 자율성 기준 명문화 (위임 체계 수립)",
        "세대 통합형 공통 업무 기준서 개발",
        "공공 미션 기반 의사결정 가이드 제작",
        "지식 관리 시스템 도입 파일럿 운영",
    ]),
    ("중기 과제", "6 ~ 12개월", C_GREEN, [
        "거버넌스 연속성 제도화 (임기 중첩, 섀도 캐비닛)",
        "의사결정 과정 공개 투명성 채널 구축",
        "조직 기억 관리 전담 체계 수립",
        "7대 영역 사후 재진단 실시 (6개월 후)",
    ]),
]

for col, (title, period, color, tasks) in enumerate(phases):
    x = Inches(0.35) + col * Inches(4.32)
    # 헤더
    add_rect(sl, x, Inches(1.15), Inches(4.0), Inches(0.65), fill=color)
    add_text(sl, title, x, Inches(1.15), Inches(4.0), Inches(0.35),
             font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, period, x, Inches(1.48), Inches(4.0), Inches(0.28),
             font_size=Pt(9.5), color=C_WHITE, align=PP_ALIGN.CENTER)
    # 태스크 박스
    add_rect(sl, x, Inches(1.85), Inches(4.0), Inches(5.4),
             fill=C_WHITE, line=C_GRAY200, line_w=Pt(0.5))
    for j, task in enumerate(tasks):
        ty = Inches(2.05) + j * Inches(1.1)
        add_rect(sl, x+Inches(0.15), ty, Inches(0.35), Inches(0.35), fill=color)
        add_text(sl, str(j+1), x+Inches(0.15), ty, Inches(0.35), Inches(0.35),
                 font_size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, "→  " + task,
                 x+Inches(0.6), ty-Inches(0.05), Inches(3.3), Inches(0.95),
                 font_size=Pt(10), color=C_DARK)

page_num(sl, 13, TOTAL)


# ────────────────────────────────────────────
# 슬라이드 14~19 : 워크숍 (1개씩)
# ────────────────────────────────────────────
for wi, ws in enumerate(WORKSHOPS):
    sl = new_slide()
    bg(sl, C_WHITE)

    clr = ws['color']
    add_rect(sl, 0, 0, W, Inches(1.6), fill=clr)

    add_text(sl, f"{ws['no']}  |  취약 영역 대응 :  {ws['target']}",
             Inches(0.5), Inches(0.1), Inches(12.5), Inches(0.45),
             font_size=Pt(10), color=RGBColor(0xff,0xff,0xff))
    add_text(sl, '\u201c' + ws['title'] + '\u201d',
             Inches(0.5), Inches(0.52), Inches(12.0), Inches(0.65),
             font_size=Pt(22), bold=True, color=C_WHITE)
    add_text(sl, ws['sub'],
             Inches(0.5), Inches(1.15), Inches(12.0), Inches(0.4),
             font_size=Pt(12), color=RGBColor(0xff,0xff,0xff))

    # 메타 칩
    meta_items = [
        ("진행 방식", "인터뷰 + 공동 작성 / 토론·설계"),
        ("소요 기간", ws['days']),
        ("추진 시점", ws['timing']),
    ]
    for mi, (mk, mv) in enumerate(meta_items):
        mx = Inches(0.4) + mi * Inches(4.3)
        add_rect(sl, mx, Inches(1.85), Inches(4.0), Inches(0.65),
                 fill=C_GRAY50, line=C_GRAY200, line_w=Pt(0.5))
        add_text(sl, mk, mx+Inches(0.12), Inches(1.88), Inches(3.8), Inches(0.28),
                 font_size=Pt(8.5), color=C_GRAY500)
        add_text(sl, mv, mx+Inches(0.12), Inches(2.13), Inches(3.8), Inches(0.32),
                 font_size=Pt(11), bold=True, color=C_DARK)

    # 워크숍 구성
    add_text(sl, "워크숍 구성",
             Inches(0.4), Inches(2.7), Inches(12.5), Inches(0.38),
             font_size=Pt(12), bold=True, color=C_NAVY)

    for di, (day_label, day_content) in enumerate([
        ("Day 1", ws['day1']), ("Day 2", ws['day2'])
    ]):
        dx = Inches(0.4) + di * Inches(6.35)
        add_rect(sl, dx, Inches(3.12), Inches(6.05), Inches(1.35),
                 fill=C_GRAY50, line=C_GRAY200, line_w=Pt(0.5))
        add_rect(sl, dx, Inches(3.12), Inches(0.9), Inches(1.35), fill=clr)
        add_text(sl, day_label, dx, Inches(3.12), Inches(0.9), Inches(1.35),
                 font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, day_content,
                 dx+Inches(1.0), Inches(3.2), Inches(4.9), Inches(1.2),
                 font_size=Pt(10.5), color=C_DARK)

    # 산출물
    add_rect(sl, Inches(0.4), Inches(4.65), Inches(12.5), Inches(1.15),
             fill=RGBColor(0xef,0xf6,0xff) if clr == C_BLUE else
                  RGBColor(0xfe,0xf2,0xf2) if clr == C_RED else
                  RGBColor(0xff,0xf7,0xed),
             line=clr, line_w=Pt(1.5))
    add_rect(sl, Inches(0.4), Inches(4.65), Inches(0.08), Inches(1.15), fill=clr)
    add_text(sl, "📋  산출물",
             Inches(0.6), Inches(4.7), Inches(3.0), Inches(0.38),
             font_size=Pt(10.5), bold=True, color=clr)
    add_text(sl, ws['output'],
             Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.6),
             font_size=Pt(10.5), color=C_DARK)

    # 기대 점수 향상
    add_text(sl, f"이 워크숍 완료 후 해당 영역 목표 점수 :  현재 → 목표",
             Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.35),
             font_size=Pt(9), color=C_GRAY500)

    page_num(sl, 14+wi, TOTAL)


# ────────────────────────────────────────────
# 슬라이드 20 : 워크숍 추진 일정 + 결론
# ────────────────────────────────────────────
sl = new_slide()
bg(sl, C_GRAY50)
nav_bar(sl, "워크숍 추진 일정 & 결론")

add_text(sl, "권장 추진 로드맵  (12개월)",
         Inches(0.5), Inches(0.65), Inches(12), Inches(0.38),
         font_size=Pt(12), bold=True, color=C_NAVY)

# 타임라인 바
timeline = [
    ("WS-01  조직 기억 복원",       [1,2], C_RED),
    ("WS-02  세대 간 기준 정립",     [3,4], C_RED),
    ("WS-03  거버넌스 연속성",       [3,4], C_ORANGE),
    ("WS-04  자율성 기준 명문화",    [5,6], C_ORANGE),
    ("WS-05  공공 미션 작동 체계",   [5,6], C_BLUE),
    ("WS-06  신뢰/공정성 재설계",    [7,12], C_BLUE),
]

month_w  = Inches(0.82)
bar_x0   = Inches(3.0)
bar_y0   = Inches(1.2)
bar_row_h= Inches(0.46)

# 월 헤더
months = ["1-2월","3-4월","5-6월","7-8월","9-10월","11-12월"]
for mi, mo in enumerate(months):
    add_text(sl, mo,
             bar_x0 + mi*month_w, bar_y0, month_w, Inches(0.32),
             font_size=Pt(8.5), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

for ri, (name, months_range, color) in enumerate(timeline):
    y = bar_y0 + Inches(0.38) + ri * bar_row_h
    add_text(sl, name, Inches(0.35), y, Inches(2.55), bar_row_h-Inches(0.05),
             font_size=Pt(9), bold=True, color=C_DARK)
    # 회색 빈 바
    add_rect(sl, bar_x0, y+Inches(0.06), month_w*6, bar_row_h-Inches(0.12),
             fill=C_GRAY100)
    # 채운 바
    s_col = (months_range[0]-1)//2
    e_col = (months_range[1]-1)//2
    add_rect(sl, bar_x0 + s_col*month_w, y+Inches(0.06),
             (e_col-s_col+1)*month_w, bar_row_h-Inches(0.12), fill=color)

# 기대 효과 박스
add_text(sl, "기대 효과",
         Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.38),
         font_size=Pt(12), bold=True, color=C_NAVY)

effects = [
    ("75점+", "목표 종합 점수\n현재 54점 → 6개월 내", C_BLUE),
    ("6개",   "실행 가능 산출물\n각 워크숍 현장 결과물", C_GREEN),
    ("100%",  "참여자 합의 기반\n구성원 스스로 만든 기준", C_ORANGE),
    ("12개월","전체 변화 기간\n단계별 구조 전환", C_NAVY),
]
for ei, (num, desc, clr) in enumerate(effects):
    ex = Inches(0.4) + ei * Inches(3.2)
    add_rect(sl, ex, Inches(5.05), Inches(3.0), Inches(1.45),
             fill=C_WHITE, line=C_GRAY200, line_w=Pt(0.5))
    add_text(sl, num, ex, Inches(5.1), Inches(3.0), Inches(0.7),
             font_size=Pt(26), bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text(sl, desc, ex, Inches(5.78), Inches(3.0), Inches(0.72),
             font_size=Pt(8.5), color=C_GRAY500, align=PP_ALIGN.CENTER)

# 결론 멘트
add_rect(sl, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.62),
         fill=C_NAVY)
add_text(sl,
         "54점은 위기가 아닌 변화의 출발점입니다.  에스비컨설팅 조직진단 AI플랫폼이 대한치과위생사협회의 지속 가능한 운영 체계를 함께 설계합니다.",
         Inches(0.6), Inches(6.72), Inches(12.2), Inches(0.58),
         font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

page_num(sl, 20, TOTAL)


# ── 저장 ────────────────────────────────────
out_path = os.path.join(
    os.path.dirname(__file__),
    "대한치과위생사협회_조직진단결과리포트.pptx"
)
prs.save(out_path)
print(f"✅ 저장 완료: {out_path}")
